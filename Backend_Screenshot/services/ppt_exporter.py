"""
PPT Report Generator
Generates PPTX matching the exact format of the provided template:
- Uses CA01437_Banoffee_English_Banner_Screenshots_Apr'26.pptx as the base
- Keeps the original cover slide
- Desktop images use Layout 0
- Mobile images use Layout 5 (two images per slide)
"""

import logging
import os
from io import BytesIO
from typing import List, Optional
from urllib.parse import urlparse

from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image

from database.db import SessionLocal
from models.screenshot import ScreenshotResult

logger = logging.getLogger(__name__)

# ── Resolve paths ─────────────────────────────────────────────────────────────
_BACKEND_ROOT = os.path.dirname(os.path.abspath(__file__))   # services/
_BACKEND_ROOT = os.path.dirname(_BACKEND_ROOT)               # Backend_Screenshot/
_PPT_FORMAT_DIR = os.path.join(_BACKEND_ROOT, "PPT_Format")

TEMPLATE_NAME = "CA01437_Banoffee_English_Banner_Screenshots_Apr'26.pptx"

# ── Low-level helpers ─────────────────────────────────────────────────────────

def _resolve_screenshot(path: str) -> Optional[str]:
    if not path:
        return None
    if os.path.isabs(path) and os.path.isfile(path):
        return path
    abs_path = os.path.join(_BACKEND_ROOT, path.replace("\\", "/").lstrip("/"))
    return abs_path if os.path.isfile(abs_path) else None

def get_fit_dimensions(img_path: str, max_w_in: float, max_h_in: float):
    try:
        with Image.open(img_path) as img:
            img_w, img_h = img.size
        # Calculate ratio to fit within max bounds
        ratio = min(max_w_in / img_w, max_h_in / img_h)
        return img_w * ratio, img_h * ratio
    except Exception as e:
        logger.error(f"Error reading image {img_path}: {e}")
        return max_w_in, max_h_in

def add_centered_picture(slide, img_path: str, box_left: float, box_top: float, box_w: float, box_h: float):
    fit_w, fit_h = get_fit_dimensions(img_path, box_w, box_h)
    
    # Center within the defined bounding box
    offset_x = (box_w - fit_w) / 2
    offset_y = (box_h - fit_h) / 2
    
    try:
        slide.shapes.add_picture(
            img_path, 
            Inches(box_left + offset_x), 
            Inches(box_top + offset_y), 
            width=Inches(fit_w), 
            height=Inches(fit_h)
        )
    except Exception as exc:
        logger.warning(f"Could not add picture to slide: {exc}")

# ── Slide Generation ──────────────────────────────────────────────────────────

def _add_desktop_slide(prs, row: ScreenshotResult):
    layout = prs.slide_layouts[0]  # Title Slide (Layout 0)
    slide = prs.slides.add_slide(layout)
    
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 10:
            # Show domain only — full URL overflows the text box
            raw_url = row.url or "—"
            try:
                domain = urlparse(raw_url).netloc or raw_url
                domain = domain.replace("www.", "")
            except Exception:
                domain = raw_url
            shape.text = f"Site: {domain}"
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(14)
        elif shape.placeholder_format.idx == 11:
            shape.text = f"Ad Size: {row.matched_creative_size or '—'}"
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(14)
        elif shape.placeholder_format.idx == 12:
            device_str = row.device or "Desktop"
            shape.text = f"Device: {device_str.capitalize()}"
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(14)
            
    ss_path = _resolve_screenshot(row.screenshot_path)
    if ss_path:
        # Based on bounds analysis for layout 0
        add_centered_picture(slide, ss_path, 0.225, 0.916, 12.905, 6.367)

def _add_mobile_slides(prs, mobile_rows: List[ScreenshotResult]):
    layout = prs.slide_layouts[5]  # Custom Layout (Layout 5)
    
    for i in range(0, len(mobile_rows), 2):
        slide = prs.slides.add_slide(layout)
        row_left = mobile_rows[i]
        row_right = mobile_rows[i+1] if i + 1 < len(mobile_rows) else None
        
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 12:
                # Left text placeholder
                device_str = row_left.device or "Mobile"
                try:
                    domain_l = urlparse(row_left.url or "").netloc.replace("www.", "") or row_left.url or "—"
                except Exception:
                    domain_l = row_left.url or "—"
                shape.text = f"Site: {domain_l}\nAd Size: {row_left.matched_creative_size or '—'}\nDevice: {device_str.capitalize()} "
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(14)
            elif shape.placeholder_format.idx == 13:
                # Right text placeholder
                if row_right:
                    device_str2 = row_right.device or "Mobile"
                    try:
                        domain_r = urlparse(row_right.url or "").netloc.replace("www.", "") or row_right.url or "—"
                    except Exception:
                        domain_r = row_right.url or "—"
                    shape.text = f"Site: {domain_r}\nAd Size: {row_right.matched_creative_size or '—'}\nDevice: {device_str2.capitalize()} "
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(14)
                else:
                    shape.text = ""

        # Add left image
        ss_left = _resolve_screenshot(row_left.screenshot_path)
        if ss_left:
            add_centered_picture(slide, ss_left, 0.281, 0.305, 3.887, 6.927)
            
        # Add right image
        if row_right:
            ss_right = _resolve_screenshot(row_right.screenshot_path)
            if ss_right:
                add_centered_picture(slide, ss_right, 9.094, 0.305, 3.880, 6.927)

# ── Public entry point ────────────────────────────────────────────────────────

def generate_ppt_report(ids: Optional[List[int]]) -> Optional[BytesIO]:
    """Generate a PPTX for the given result IDs. Pass None to include all results."""
    db = SessionLocal()
    try:
        q = db.query(ScreenshotResult)
        if ids is not None:
            q = q.filter(ScreenshotResult.id.in_(ids))
        rows = q.filter(ScreenshotResult.status == "success") \
                .order_by(ScreenshotResult.created_at) \
                .all()
    finally:
        db.close()

    if not rows:
        logger.warning("generate_ppt_report: no rows found for ids=%s", ids)
        return None

    template_path = os.path.join(_PPT_FORMAT_DIR, TEMPLATE_NAME)
    if not os.path.exists(template_path):
        logger.error(f"Template not found at {template_path}")
        return None

    prs = Presentation(template_path)

    # Keep ONLY the cover slide, delete the rest of the existing example slides safely
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    for i in range(len(slides) - 1, 0, -1):
        # Drop the relationship to avoid PPT corruption warnings
        rId = slides[i].rId
        prs.part.drop_rel(rId)
        xml_slides.remove(slides[i])

    # Separate rows by device
    desktop_rows = []
    mobile_rows = []
    for r in rows:
        device_lower = str(r.device).lower() if r.device else "desktop"
        if device_lower == "mobile":
            mobile_rows.append(r)
        else:
            desktop_rows.append(r)

    # Add Desktop slides (1 image per slide)
    for row in desktop_rows:
        _add_desktop_slide(prs, row)
        
    # Add Mobile slides (2 images per slide)
    if mobile_rows:
        _add_mobile_slides(prs, mobile_rows)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    logger.info("PPT report generated: %d desktop, %d mobile results", len(desktop_rows), len(mobile_rows))
    return buf
